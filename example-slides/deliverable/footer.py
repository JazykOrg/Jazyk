"""The footer applied to every slide: site + copyright left, author right."""

from pptx.enum.text import PP_ALIGN
from pptx.util import Inches

import theme

SITE = "https://jazyk.org"
COPYRIGHT = "© Jazyk"
AUTHOR = "Matus Faro"


def add_footer(slide, prs):
    """Add the template footer to one slide."""
    top = prs.slide_height - Inches(0.6)

    left_box = slide.shapes.add_textbox(Inches(0.4), top, Inches(6), Inches(0.4))
    left_para = left_box.text_frame.paragraphs[0]
    left_run = left_para.add_run()
    left_run.text = "%s %s" % (SITE, COPYRIGHT)
    left_run.hyperlink.address = SITE
    theme.style_run(left_run, size=theme.FOOTER_SIZE)

    right_box = slide.shapes.add_textbox(
        prs.slide_width - Inches(6.4), top, Inches(6), Inches(0.4)
    )
    right_para = right_box.text_frame.paragraphs[0]
    right_para.alignment = PP_ALIGN.RIGHT
    right_run = right_para.add_run()
    right_run.text = AUTHOR
    theme.style_run(right_run, size=theme.FOOTER_SIZE)
