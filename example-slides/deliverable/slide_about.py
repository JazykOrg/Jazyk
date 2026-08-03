"""The About slide: what Jazyk is about, in a couple of sentences."""

from pptx.util import Inches, Pt

import footer
import theme

ABOUT_TEXT = (
    "Jazyk is a natural language compiler: it treats plain prose documentation "
    "as the source code of a program. From that prose it maintains a semantic "
    "graph of entities and requirements, which downstream tools consume to "
    "generate code, tests, and project plans."
)


def build(prs):
    """Add the About slide to the presentation."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank layout

    title_box = slide.shapes.add_textbox(Inches(1), Inches(0.7), Inches(11.333), Inches(1.0))
    title_para = title_box.text_frame.paragraphs[0]
    title_run = title_para.add_run()
    title_run.text = "About Jazyk"
    theme.style_run(title_run, size=Pt(40), color=theme.PRIMARY, bold=True)

    # Accent bar in the complementary color, per the template palette.
    bar = slide.shapes.add_shape(1, Inches(1), Inches(1.8), Inches(2.5), Inches(0.08))
    bar.fill.solid()
    bar.fill.fore_color.rgb = theme.COMPLEMENTARY
    bar.line.fill.background()

    body_box = slide.shapes.add_textbox(Inches(1), Inches(2.2), Inches(11.333), Inches(3.5))
    body_box.text_frame.word_wrap = True
    body_para = body_box.text_frame.paragraphs[0]
    body_run = body_para.add_run()
    body_run.text = ABOUT_TEXT
    theme.style_run(body_run, size=theme.BODY_SIZE)

    footer.add_footer(slide, prs)
    return slide
