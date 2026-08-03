"""The Introduction slide: headline title and site link."""

from pptx.util import Inches, Pt

import footer
import theme

SITE = "https://jazyk.org"


def build(prs):
    """Add the Introduction slide to the presentation."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank layout

    title_box = slide.shapes.add_textbox(Inches(1), Inches(2.4), Inches(11.333), Inches(1.6))
    title_para = title_box.text_frame.paragraphs[0]
    title_run = title_para.add_run()
    title_run.text = "Jazyk"
    theme.style_run(title_run, size=theme.TITLE_SIZE, color=theme.PRIMARY, bold=True)

    link_box = slide.shapes.add_textbox(Inches(1), Inches(4.2), Inches(11.333), Inches(0.8))
    link_para = link_box.text_frame.paragraphs[0]
    link_run = link_para.add_run()
    link_run.text = SITE
    link_run.hyperlink.address = SITE
    theme.style_run(link_run, size=Pt(24), color=theme.COMPLEMENTARY)

    footer.add_footer(slide, prs)
    return slide
