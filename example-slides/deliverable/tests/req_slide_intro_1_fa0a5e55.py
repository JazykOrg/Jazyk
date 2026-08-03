"""The Introduction slide shall show a headline title `Jazyk`."""

import sys

from pptx import Presentation


def req_slide_intro_1_fa0a5e55():
    prs = Presentation("dist/jazyk-slides.pptx")
    intro = list(prs.slides)[0]
    headline = None
    for shape in intro.shapes:
        if shape.has_text_frame and shape.text_frame.text.strip() == "Jazyk":
            headline = shape
            break
    assert headline is not None, "Introduction slide has no `Jazyk` headline"
    run = headline.text_frame.paragraphs[0].runs[0]
    size = run.font.size
    assert size is not None and size.pt >= 40, (
        "the `Jazyk` title is not headline-sized: %r" % size
    )


if __name__ == "__main__":
    req_slide_intro_1_fa0a5e55()
    print("PASS req:slide-intro-1")
    sys.exit(0)
