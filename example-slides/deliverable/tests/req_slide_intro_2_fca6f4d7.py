"""The Introduction slide shall show a link to the site https://jazyk.org."""

import sys

from pptx import Presentation

SITE = "https://jazyk.org"


def req_slide_intro_2_fca6f4d7():
    prs = Presentation("dist/jazyk-slides.pptx")
    intro = list(prs.slides)[0]
    found = False
    for shape in intro.shapes:
        if not shape.has_text_frame:
            continue
        for para in shape.text_frame.paragraphs:
            for run in para.runs:
                if run.hyperlink.address == SITE and SITE in run.text:
                    found = True
    assert found, "Introduction slide has no visible hyperlink to %s" % SITE


if __name__ == "__main__":
    req_slide_intro_2_fca6f4d7()
    print("PASS req:slide-intro-2")
    sys.exit(0)
