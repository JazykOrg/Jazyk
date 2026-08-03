"""The slide deck shall include the Introduction slide as its first slide."""

import sys

from pptx import Presentation


def req_slides_1_24d654ff():
    prs = Presentation("dist/jazyk-slides.pptx")
    slides = list(prs.slides)
    assert len(slides) >= 1, "deck has no slides"
    first_texts = [sh.text_frame.text for sh in slides[0].shapes if sh.has_text_frame]
    assert any(t.strip() == "Jazyk" for t in first_texts), (
        "first slide is not the Introduction slide (no `Jazyk` headline); texts: %r" % first_texts
    )


if __name__ == "__main__":
    req_slides_1_24d654ff()
    print("PASS req:slides-1")
    sys.exit(0)
