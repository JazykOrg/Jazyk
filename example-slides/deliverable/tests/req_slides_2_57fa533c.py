"""The slide deck shall include the About slide as its second slide."""

import sys

from pptx import Presentation


def req_slides_2_57fa533c():
    prs = Presentation("dist/jazyk-slides.pptx")
    slides = list(prs.slides)
    assert len(slides) >= 2, "deck has fewer than two slides"
    second_texts = [sh.text_frame.text for sh in slides[1].shapes if sh.has_text_frame]
    assert any("About" in t for t in second_texts), (
        "second slide is not the About slide; texts: %r" % second_texts
    )


if __name__ == "__main__":
    req_slides_2_57fa533c()
    print("PASS req:slides-2")
    sys.exit(0)
