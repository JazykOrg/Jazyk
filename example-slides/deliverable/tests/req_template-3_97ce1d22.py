import unittest
from pptx import Presentation
from pptx.dml.color import RGBColor

class TestTemplate3(unittest.TestCase):
    def setUp(self):
        self.filepath = "dist/jazyk-slides.pptx"
        try:
            prs = Presentation(self.filepath)
        except FileNotFoundError:
            self.fail("The artifact dist/jazyk-slides.pptx was not found after the build.")

    def test_primary_color_is_correct(self):
        # Check the color used in the accent bar of the About slide (Slide 2).
        # The accent bar is added in slide_about.py and uses theme.COMPLEMENTARY, but let's check a title element if possible or rely on the general template usage.
        # Since the requirement specifies #248555 for primary color, we should check elements styled with PRIMARY.
        # The Introduction slide title is styled with PRIMARY.
        slide_intro = prs.slides[0]
        title_shape = slide_intro.shapes[0]
        
        # Assuming the first run in the title box (which is the main title) has the correct color.
        if title_shape.has_text_frame:
            first_paragraph = title_shape.text_frame.paragraphs[0]
            title_run = first_paragraph.runs[0]
            
            # Check if the run's font color is set to #248555 (RGBColor.from_string("248555"))
            self.assertEqual(title_run.font.color.rgb, RGBColor.from_string("248555"), "The primary color of the title must be #248555.")

if __name__ == '__main__':
    unittest.main()