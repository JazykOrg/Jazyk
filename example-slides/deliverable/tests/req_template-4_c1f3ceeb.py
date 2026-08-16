import unittest
from pptx import Presentation
from pptx.dml.color import RGBColor

class TestTemplate4(unittest.TestCase):
    def setUp(self):
        self.filepath = "dist/jazyk-slides.pptx"
        try:
            prs = Presentation(self.filepath)
        except FileNotFoundError:
            self.fail("The artifact dist/jazyk-slides.pptx was not found after the build.")

    def test_complementary_color_is_correct(self):
        # Check the color used for the accent bar on the About slide (Slide 2).
        slide_about = prs.slides[1]
        
        # The accent bar is added in slide_about.py and uses theme.COMPLEMENTARY.
        # We need to find the shape that represents the accent bar. It's usually a simple rectangle shape (type 1).
        # Since we don't know the exact order of shapes, let's look for a shape added near the top with dimensions matching the bar.
        
        accent_bar = None
        for shape in slide_about.shapes:
            if shape.shape_type == 1 and shape.top < 2.0 and shape.height < 0.1: # Rough check for a small horizontal bar near the top
                accent_bar = shape
                break

        self.assertIsNotNone(accent_bar, "Accent bar shape was not found on the About slide.")
        
        # Check the fill color of the accent bar
        fill_color = accent_bar.fill.fore_color.rgb
        expected_color = RGBColor.from_string("8E4367")
        self.assertEqual(fill_color, expected_color, "The complementary color of the accent bar must be #8e4367.")

if __name__ == '__main__':
    unittest.main()