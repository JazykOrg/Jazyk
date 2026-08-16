import unittest
from pptx import Presentation

class TestTemplate5(unittest.TestCase):
    def setUp(self):
        self.filepath = "dist/jazyk-slides.pptx"
        try:
            prs = Presentation(self.filepath)
        except FileNotFoundError:
            self.fail("The artifact dist/jazyk-slides.pptx was not found after the build.")

    def test_font_is_comic_sans_ms(self):
        # Check a body text element, e.g., in the About slide (Slide 2).
        slide_about = prs.slides[1]
        body_shape = None
        for shape in slide_about.shapes:
            if shape.has_text_frame and shape.top > 2.0: # Rough check for the body text box
                body_shape = shape
                break

        self.assertIsNotNone(body_shape, "Body text shape was not found on the About slide.")
        
        # Check the font name of the run containing the main description text.
        if body_shape.has_text_frame:
            first_paragraph = body_shape.text_frame.paragraphs[0]
            body_run = first_paragraph.runs[0]
            self.assertEqual(body_run.font.name, "Comic Sans MS", "The font used for the body text must be Comic Sans MS.")

if __name__ == '__main__':
    unittest.main()