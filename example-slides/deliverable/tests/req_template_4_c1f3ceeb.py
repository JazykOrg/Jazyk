"""The slide deck shall use #8e4367 as the complementary color."""

import sys

from pptx import Presentation
from pptx.dml.color import RGBColor

TARGET = RGBColor.from_string("8E4367")


def _colors(prs):
    for slide in prs.slides:
        for shape in slide.shapes:
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    for run in para.runs:
                        c = run.font.color
                        if c and c.type is not None and c.rgb is not None:
                            yield c.rgb
            try:
                fill = shape.fill
                if fill.type is not None and fill.fore_color.rgb is not None:
                    yield fill.fore_color.rgb
            except (TypeError, AttributeError):
                pass


def req_template_4_c1f3ceeb():
    prs = Presentation("dist/jazyk-slides.pptx")
    used = set(_colors(prs))
    assert TARGET in used, "complementary color #8e4367 not used anywhere; colors used: %r" % used


if __name__ == "__main__":
    req_template_4_c1f3ceeb()
    print("PASS req:template-4")
    sys.exit(0)
