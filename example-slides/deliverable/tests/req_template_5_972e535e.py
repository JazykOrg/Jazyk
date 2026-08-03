"""The slide deck shall use the Comis Sans font."""

import sys

from pptx import Presentation

FONT = "Comis Sans"


def req_template_5_972e535e():
    prs = Presentation("dist/jazyk-slides.pptx")
    checked = 0
    for slide in prs.slides:
        for shape in slide.shapes:
            if not shape.has_text_frame:
                continue
            for para in shape.text_frame.paragraphs:
                for run in para.runs:
                    if not run.text.strip():
                        continue
                    checked += 1
                    assert run.font.name == FONT, (
                        "run %r uses font %r, expected %r" % (run.text, run.font.name, FONT)
                    )
    assert checked > 0, "deck has no text runs to check"


if __name__ == "__main__":
    req_template_5_972e535e()
    print("PASS req:template-5")
    sys.exit(0)
