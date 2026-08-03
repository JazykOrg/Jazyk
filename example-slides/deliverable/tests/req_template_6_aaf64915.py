"""The footer shall show the site https://jazyk.org with a copyright at the bottom left."""

import sys

from pptx import Presentation
from pptx.util import Emu

SITE = "https://jazyk.org"


def req_template_6_aaf64915():
    prs = Presentation("dist/jazyk-slides.pptx")
    mid_x = prs.slide_width // 2
    mid_y = prs.slide_height // 2
    for idx, slide in enumerate(prs.slides):
        found = False
        for shape in slide.shapes:
            if not shape.has_text_frame:
                continue
            text = shape.text_frame.text
            if SITE in text and "©" in text:
                assert shape.left < mid_x, "slide %d: site+copyright is not on the left" % idx
                assert shape.top > mid_y, "slide %d: site+copyright is not at the bottom" % idx
                found = True
        assert found, "slide %d has no bottom-left site+copyright footer text" % idx


if __name__ == "__main__":
    req_template_6_aaf64915()
    print("PASS req:template-6")
    sys.exit(0)
