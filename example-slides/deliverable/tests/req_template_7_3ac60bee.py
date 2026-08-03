"""The footer shall show the author Matus Faro at the bottom right."""

import sys

from pptx import Presentation

AUTHOR = "Matus Faro"


def req_template_7_3ac60bee():
    prs = Presentation("dist/jazyk-slides.pptx")
    mid_x = prs.slide_width // 2
    mid_y = prs.slide_height // 2
    for idx, slide in enumerate(prs.slides):
        found = False
        for shape in slide.shapes:
            if not shape.has_text_frame:
                continue
            if AUTHOR in shape.text_frame.text:
                right_edge = shape.left + shape.width
                assert right_edge > mid_x, "slide %d: author box is not on the right" % idx
                assert shape.top > mid_y, "slide %d: author box is not at the bottom" % idx
                found = True
        assert found, "slide %d has no bottom-right author footer text" % idx


if __name__ == "__main__":
    req_template_7_3ac60bee()
    print("PASS req:template-7")
    sys.exit(0)
