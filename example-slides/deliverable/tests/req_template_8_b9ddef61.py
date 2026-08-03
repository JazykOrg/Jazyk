"""Every slide in the slide deck shall have a footer."""

import sys

from pptx import Presentation

SITE = "https://jazyk.org"
AUTHOR = "Matus Faro"


def req_template_8_b9ddef61():
    prs = Presentation("dist/jazyk-slides.pptx")
    slides = list(prs.slides)
    assert slides, "deck has no slides"
    mid_y = prs.slide_height // 2
    for idx, slide in enumerate(slides):
        footer_texts = [
            shape.text_frame.text
            for shape in slide.shapes
            if shape.has_text_frame and shape.top is not None and shape.top > mid_y
        ]
        has_site = any(SITE in t for t in footer_texts)
        has_author = any(AUTHOR in t for t in footer_texts)
        assert has_site and has_author, (
            "slide %d lacks a complete footer (site: %s, author: %s)" % (idx, has_site, has_author)
        )


if __name__ == "__main__":
    req_template_8_b9ddef61()
    print("PASS req:template-8")
    sys.exit(0)
